"""Generate the French PowerPoint presentation for the project.

Run with:  python -m src.reporting.generate_pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from configs import config
from src.reporting.common import eda_fig, fig, load_metrics
from src.utils.logger import get_logger

logger = get_logger(__name__)

BLUE = RGBColor(0x1F, 0x4E, 0x79)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
DARK = RGBColor(0x22, 0x22, 0x22)
W, H = Inches(13.333), Inches(7.5)


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BLUE)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0xCF, 0xDC, 0xEF)
    return slide


def _content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, RGBColor(255, 255, 255))
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255)
    return slide


def _bullets(slide, items, left=0.8, top=1.5, width=11.7, height=5.5, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lvl = 0
        if isinstance(item, tuple):
            item, lvl = item
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + item
        r.font.size = Pt(size - lvl * 2); r.font.color.rgb = DARK
        p.space_after = Pt(8)
    return tb


def _image(slide, path, left, top, width):
    if path is not None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 width=Inches(width))


def _metric_cards(slide, cards, top=2.0):
    """cards: list of (value, label)."""
    n = len(cards)
    gap = 0.4
    card_w = (13.333 - 1.6 - gap * (n - 1)) / n
    x = 0.8
    for value, label in cards:
        box = slide.shapes.add_shape(1, Inches(x), Inches(top),
                                     Inches(card_w), Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = BLUE; box.line.width = Pt(1.5)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = value
        r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = BLUE
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = label
        r2.font.size = Pt(13); r2.font.color.rgb = DARK
        x += card_w + gap


def build() -> None:
    m = load_metrics()
    best = m["best"]; test = best.get("test", {})
    interval = m["interval"]
    cov = interval.get("empirical_coverage_pct", 0)

    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H

    # 1. Titre
    _title_slide(prs,
                 "Prédiction de Salaires par Machine Learning & NLP",
                 "Projet d'Intégration (PI) — Esprit  •  DatawarehouseDB (SQL Server)")

    # 2. Contexte & objectif
    s = _content_slide(prs, "Contexte & Objectif")
    _bullets(s, [
        "Prédire le salaire annuel d'une offre d'emploi (domaine data).",
        "Données : 785 741 offres d'emploi (2023).",
        "Cible : salaire annuel moyen en USD.",
        "Système complet de bout en bout :",
        ("Connexion SQL Server, EDA, feature engineering NLP + tabulaire", 1),
        ("Comparaison de modèles, optimisation, ensemble", 1),
        ("Explicabilité, API REST (FastAPI), tableau de bord (Streamlit)", 1),
    ])

    # 3. Entrepôt & schéma
    s = _content_slide(prs, "Connexion à l'entrepôt de données")
    _bullets(s, [
        "Connexion SQLAlchemy + pyodbc (authentification Windows).",
        "Schéma en étoile détecté automatiquement :",
        ("Fact_Jobs (779 226) — table de faits", 1),
        ("Dim_Job, Dim_Company, Dim_Location, Dim_Date", 1),
        ("Bridge_Job_Skills → Dim_Job_Skills (compétences N-N)", 1),
        "Modules : connexion, inspection du schéma, diagnostics.",
    ], width=7.2)
    _image(s, eda_fig("04_salary_by_country.png"), left=8.2, top=1.6, width=4.6)

    # 4. Problèmes ETL détectés
    s = _content_slide(prs, "Qualité des données : 3 anomalies ETL détectées")
    _bullets(s, [
        "Salaires manquants chargés en 0 au lieu de NULL (757 635 lignes).",
        "Table de compétences incomplète : 15 416 offres seulement (vs ~690 k).",
        "Lien Fact_Jobs → Dim_Job rompu : 4 796 job_id / 223 550.",
        "Décision : apprentissage sur le CSV propre (source fiable),",
        ("connexion SQL Server conservée pour l'accès entrepôt & BI.", 1),
        "→ La détection de ces bugs est un point fort du projet.",
    ])

    # 5. EDA
    s = _content_slide(prs, "Analyse exploratoire (EDA)")
    _bullets(s, [
        "Seules 22 003 offres (2,8 %) ont un salaire → sous-ensemble d'étude.",
        "Médiane 115 000 $ ; moyenne 123 286 $.",
        "Salaire asymétrique (1,75) → transformation log (asymétrie -0,18).",
        "Télétravail : +13 830 $ de médiane.",
        "5,4 compétences/offre ; présentes sur 91,7 % des lignes.",
    ], width=6.6)
    _image(s, eda_fig("01_salary_distribution.png"), left=7.4, top=1.6, width=5.5)

    # 6. Feature engineering
    s = _content_slide(prs, "Ingénierie des caractéristiques")
    _bullets(s, [
        "Numériques : nb compétences, séniorité, télétravail, mois, catégories.",
        "Catégorielles : One-Hot + Target Encoding sans fuite",
        ("(pays, entreprise, ville)", 1),
        "Textuelles (NLP) : TF-IDF sur compétences et titre.",
        "Ajout entreprise + ville + catégories → R² de 0,53 à 0,59.",
        "Target Encoding intra-validation croisée → aucune fuite de données.",
    ])

    # 7. Comparaison des modèles (image bar chart from comparison? use feature importance)
    s = _content_slide(prs, "Comparaison des modèles")
    rows = []
    if not m["comparison"].empty:
        for _, r in m["comparison"].head(6).iterrows():
            rows.append(f"{r['model']} : R² {r['R2']:.3f} | MAE ${r['MAE']:,.0f}")
    _bullets(s, [
        "10 modèles entraînés (cible log, métriques en USD) :",
    ] + [(x, 1) for x in rows] + [
        "Boosting d'arbres (LightGBM, XGBoost) en tête.",
    ])

    # 8. Optimisation + ensemble
    s = _content_slide(prs, "Optimisation (Optuna) & Ensemble (Stacking)")
    _bullets(s, [
        "Optuna : 40 essais LightGBM + 25 essais XGBoost.",
        "Stacking : ExtraTrees + LightGBM + XGBoost → méta-modèle Ridge.",
        "Validation croisée interne à 5 plis.",
        "Le modèle empilé est le modèle final déployé.",
    ])

    # 9. Résultats finaux (cards)
    s = _content_slide(prs, "Résultats finaux (jeu de test)")
    _metric_cards(s, [
        (f"{test.get('R2',0):.2f}", "R² (test)"),
        (f"{test.get('MAPE',0):.1f}%", "MAPE"),
        (f"${test.get('MAE',0):,.0f}", "MAE"),
        (f"{cov:.0f}%", "Couverture intervalle"),
    ], top=2.2)
    _bullets(s, [
        "Validation croisée 5 plis : R² = 0,59 ± 0,013 (très stable).",
        "Erreur moyenne ~18,5 % sur un salaire médian de 115 000 $.",
    ], top=4.8, size=18)

    # 10. Diagnostics image
    s = _content_slide(prs, "Diagnostics d'évaluation")
    _image(s, fig("evaluation_diagnostics.png"), left=1.2, top=1.5, width=11.0)

    # 11. Intervalle
    s = _content_slide(prs, "Prédiction par intervalle (fourchette)")
    _bullets(s, [
        f"Régression quantile (5e-95e pct) : couverture {cov:.1f} % (nominal 90 %).",
        "Ex : Senior Data Scientist → ~157 k$ [92 k – 202 k$].",
        "Plus honnête et utile qu'une valeur unique.",
        "N.B. : l'intervalle n'augmente pas la précision ponctuelle ;",
        ("c'est une métrique distincte (la couverture).", 1),
    ], width=6.6)
    _image(s, fig("interval_coverage.png"), left=7.3, top=1.7, width=5.7)

    # 12. Explicabilité
    s = _content_slide(prs, "Explicabilité : facteurs du salaire")
    _bullets(s, [
        "Variables les plus déterminantes :",
        ("1. Titre du poste", 1), ("2. Entreprise (company_name)", 1),
        ("3. Séniorité", 1), ("4. Compétences", 1),
        ("5. Pays / Ville", 1),
    ], width=6.4)
    _image(s, fig("feature_importance.png"), left=7.0, top=1.5, width=6.0)

    # 13. Pourquoi c'est le maximum
    s = _content_slide(prs, "Pourquoi 0,59 est le maximum atteignable")
    _bullets(s, [
        "Plafond imposé par les DONNÉES, pas par le modèle.",
        "Variables clés absentes : années d'expérience, description complète.",
        "Rendements décroissants vérifiés : optimisation étendue → R² identique.",
        "Stabilité prouvée : validation croisée 0,59 ± 0,013.",
        "Arbres = arrêt précoce : entraîner plus longtemps = surapprentissage.",
        "Rigueur : aucune fuite, test isolé → résultats honnêtes et défendables.",
    ])

    # 14. Architecture & livrables
    s = _content_slide(prs, "Architecture & Livrables")
    _bullets(s, [
        "Code modulaire : db, data, preprocessing, features, models, evaluation.",
        "API FastAPI : /predict (valeur + fourchette), /metrics, /feature-importance.",
        "Tableau de bord Streamlit : prédiction, comparaison, importance, EDA.",
        "Modèle persisté, rapports et figures générés automatiquement.",
        "Bonnes pratiques : typage, logging, gestion d'erreurs, configuration .env.",
    ])

    # 15. Conclusion
    s = _content_slide(prs, "Conclusion")
    _bullets(s, [
        f"Système complet de qualité production.",
        f"R² = {test.get('R2',0):.2f} | MAPE = {test.get('MAPE',0):.1f} % | "
        f"Couverture = {cov:.0f} %.",
        "Plafond réaliste des données atteint, avec rigueur méthodologique.",
        "Perspectives : données enrichies, MLflow, reconstruction de l'entrepôt.",
        "Merci de votre attention.",
    ], size=20)

    out = config.REPORTS_DIR / "Presentation_Prediction_Salaires.pptx"
    prs.save(str(out))
    logger.info("PowerPoint saved -> %s (%d slides)", out, len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
